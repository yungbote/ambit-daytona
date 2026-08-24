from __future__ import annotations

import argparse
import datetime as dt
import json
import sys
import zipfile
from pathlib import Path

from openpyxl import Workbook, load_workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.styles import Font, PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.workbook.defined_name import DefinedName
from openpyxl.worksheet.table import Table, TableStyleInfo
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from common import canonical_json, file_receipts, runtime_guard, sha256


FIXED_TIME = dt.datetime(2026, 8, 23, 12, 0, 0)
PACK_REF = "ambit.runtime-pack/office-authoring@1"


def canonicalize_zip(path: Path) -> None:
    with zipfile.ZipFile(path, "r") as source:
        members = [(info.filename, source.read(info.filename)) for info in source.infolist()]
    temporary = path.with_suffix(path.suffix + ".canonical")
    with zipfile.ZipFile(temporary, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9) as output:
        for name, payload in sorted(members):
            info = zipfile.ZipInfo(name, date_time=(1980, 1, 1, 0, 0, 0))
            info.compress_type = zipfile.ZIP_DEFLATED
            info.create_system = 3
            info.external_attr = 0o100444 << 16
            output.writestr(info, payload)
    temporary.replace(path)


def create_workbook(path: Path) -> None:
    workbook = Workbook()
    workbook.properties.title = "Quarterly revenue model"
    workbook.properties.subject = "Accessible deterministic spreadsheet conformance"
    workbook.properties.creator = "Ambit C18 conformance"
    workbook.properties.created = FIXED_TIME
    workbook.properties.modified = FIXED_TIME
    sheet = workbook.active
    sheet.title = "Revenue"
    sheet.sheet_properties.pageSetUpPr.fitToPage = True
    rows = [
        ("Region", "Q1", "Q2", "Total"),
        ("North", 120, 130, "=SUM(B2:C2)"),
        ("South", 90, 110, "=SUM(B3:C3)"),
        ("West", 80, 95, "=SUM(B4:C4)"),
        ("Grand total", "=SUM(B2:B4)", "=SUM(C2:C4)", "=SUM(D2:D4)"),
    ]
    for row in rows:
        sheet.append(row)
    for cell in sheet[1]:
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="1F4E78")
    for cell in sheet[5]:
        cell.font = Font(bold=True)
    for column in range(1, 5):
        sheet.column_dimensions[get_column_letter(column)].width = 18
    table = Table(displayName="RevenueTable", ref="A1:D4")
    table.tableStyleInfo = TableStyleInfo(
        name="TableStyleMedium2",
        showFirstColumn=False,
        showLastColumn=False,
        showRowStripes=True,
        showColumnStripes=False,
    )
    sheet.add_table(table)
    chart = BarChart()
    chart.title = "Revenue by region"
    chart.y_axis.title = "Amount"
    chart.x_axis.title = "Region"
    chart.add_data(Reference(sheet, min_col=4, min_row=1, max_row=4), titles_from_data=True)
    chart.set_categories(Reference(sheet, min_col=1, min_row=2, max_row=4))
    chart.height = 7
    chart.width = 12
    sheet.add_chart(chart, "F2")
    workbook.defined_names.add(DefinedName("RevenueTotals", attr_text="'Revenue'!$D$2:$D$4"))
    workbook.calculation.fullCalcOnLoad = True
    workbook.calculation.forceFullCalc = True
    workbook.calculation.calcMode = "auto"
    workbook.save(path)
    canonicalize_zip(path)


def revise_workbook(source: Path, target: Path) -> None:
    workbook = load_workbook(source)
    sheet = workbook["Revenue"]
    sheet["C2"] = 135
    workbook.properties.modified = FIXED_TIME
    workbook.save(target)
    canonicalize_zip(target)


def _set_description(shape: object, description: str) -> None:
    elements = shape._element.xpath(".//*[local-name()='cNvPr']")  # type: ignore[attr-defined]
    if not elements:
        raise ValueError("shape has no non-visual properties")
    elements[0].set("descr", description)


def create_presentation(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    presentation.core_properties.title = "Quarterly revenue briefing"
    presentation.core_properties.subject = "Accessible deterministic presentation conformance"
    presentation.core_properties.author = "Ambit C18 conformance"
    presentation.core_properties.created = FIXED_TIME
    presentation.core_properties.modified = FIXED_TIME

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "Quarterly revenue"
    title_slide.placeholders[1].text = "A deterministic C18 presentation fixture"
    title_notes = title_slide.notes_slide.notes_text_frame
    if title_notes is None:
        raise ValueError("title slide has no notes text frame")
    title_notes.text = "Introduce the source and explain the units."

    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Revenue by region"
    chart_data = ChartData()
    chart_data.categories = ["North", "South", "West"]
    chart_data.add_series("Q1", (120, 90, 80))
    chart_data.add_series("Q2", (130, 110, 95))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8),
        Inches(1.5),
        Inches(8.2),
        Inches(4.8),
        chart_data,
    )
    chart.chart.has_legend = True
    chart.chart.has_title = True
    chart.chart.chart_title.text_frame.text = "Regional revenue"
    _set_description(chart, "Column chart comparing Q1 and Q2 revenue by region")
    summary = slide.shapes.add_textbox(Inches(9.35), Inches(1.7), Inches(3.1), Inches(3.2))
    summary.text_frame.text = "Key finding"
    summary.text_frame.paragraphs[0].font.bold = True
    summary.text_frame.paragraphs[0].font.size = Pt(24)
    paragraph = summary.text_frame.add_paragraph()
    paragraph.text = "All regions grew; North remains the largest contributor."
    paragraph.font.size = Pt(18)
    paragraph.alignment = PP_ALIGN.LEFT
    _set_description(summary, "Text summary of the chart's key finding")
    notes = slide.notes_slide.notes_text_frame
    if notes is None:
        raise ValueError("chart slide has no notes text frame")
    notes.text = "Call out the revised North Q2 value when present."
    presentation.save(path)
    canonicalize_zip(path)


def revise_presentation(source: Path, target: Path) -> None:
    presentation = Presentation(source)
    presentation.slides[1].shapes.title.text = "Revenue by region — revised"
    presentation.core_properties.modified = FIXED_TIME
    presentation.save(target)
    canonicalize_zip(target)


def _validate_workbook(path: Path, *, revised: bool) -> dict[str, object]:
    workbook = load_workbook(path, data_only=False)
    sheet = workbook["Revenue"]
    assert sheet["D2"].value == "=SUM(B2:C2)"
    assert sheet["C2"].value == (135 if revised else 130)
    assert sheet.tables["RevenueTable"].ref == "A1:D4"
    assert len(sheet._charts) == 1
    assert "RevenueTotals" in workbook.defined_names
    assert sheet["A1"].style_id != 0
    assert workbook.calculation.calcMode == "auto"
    return {
        "sheetCount": len(workbook.sheetnames),
        "formulaCount": sum(
            1
            for row in sheet.iter_rows()
            for cell in row
            if isinstance(cell.value, str) and cell.value.startswith("=")
        ),
        "tableCount": len(sheet.tables),
        "chartCount": len(sheet._charts),
        "definedNameCount": len(workbook.defined_names),
    }


def _overlap(left: object, right: object) -> bool:
    return not (
        left.left + left.width <= right.left
        or right.left + right.width <= left.left
        or left.top + left.height <= right.top
        or right.top + right.height <= left.top
    )


def _validate_presentation(path: Path, *, revised: bool) -> dict[str, object]:
    presentation = Presentation(path)
    assert len(presentation.slides) == 2
    assert presentation.slides[0].shapes.title.text == "Quarterly revenue"
    slide = presentation.slides[1]
    assert slide.shapes.title.text.endswith("revised") is revised
    charts = [shape for shape in slide.shapes if getattr(shape, "has_chart", False)]
    assert len(charts) == 1
    notes_frames = [item.notes_slide.notes_text_frame for item in presentation.slides]
    assert all(frame is not None for frame in notes_frames)
    notes = [frame.text for frame in notes_frames if frame is not None]
    assert all(value.strip() for value in notes)
    title_shape_id = slide.shapes.title.shape_id
    content_shapes = [shape for shape in slide.shapes if shape.shape_id != title_shape_id]
    for shape in content_shapes:
        assert shape.left >= 0 and shape.top >= 0
        assert shape.left + shape.width <= presentation.slide_width
        assert shape.top + shape.height <= presentation.slide_height
        descriptions = shape._element.xpath(".//*[local-name()='cNvPr']/@descr")
        assert descriptions and descriptions[0]
    assert not _overlap(content_shapes[0], content_shapes[1])
    return {
        "slideCount": len(presentation.slides),
        "chartCount": len(charts),
        "notesCount": len(notes),
        "describedContentShapeCount": len(content_shapes),
        "layoutCollisionCount": 0,
    }


def generate(output: Path) -> None:
    output.mkdir(parents=True, exist_ok=False)
    spreadsheet_v1 = output / "spreadsheet-v1.xlsx"
    spreadsheet_v2 = output / "spreadsheet-v2.xlsx"
    presentation_v1 = output / "presentation-v1.pptx"
    presentation_v2 = output / "presentation-v2.pptx"
    create_workbook(spreadsheet_v1)
    revise_workbook(spreadsheet_v1, spreadsheet_v2)
    create_presentation(presentation_v1)
    revise_presentation(presentation_v1, presentation_v2)
    structures: dict[str, object] = {}
    for path in (spreadsheet_v1, spreadsheet_v2, presentation_v1, presentation_v2):
        name = path.name
        if path.suffix == ".xlsx":
            structures[name] = _validate_workbook(path, revised="v2" in name)
        else:
            structures[name] = _validate_presentation(path, revised="v2" in name)
    assert structures["spreadsheet-v1.xlsx"] == structures["spreadsheet-v2.xlsx"]
    assert structures["presentation-v1.pptx"] == structures["presentation-v2.pptx"]
    canonical_json(
        output / "fixture-manifest.json",
        {
            "schema": "ambit.c18-office-fixture-manifest/v1",
            "packRef": PACK_REF,
            "structures": structures,
            "files": file_receipts(output),
        },
    )


def finalize(output: Path) -> None:
    fixtures = output / "fixtures"
    rendered = output / "rendered"
    expected_text = {
        "spreadsheet-v1.txt": ["North", "South", "West", "250", "200", "175"],
        "spreadsheet-v2.txt": ["North", "South", "West", "255", "200", "175"],
        "presentation-v1.txt": ["Quarterly revenue", "All regions grew"],
        "presentation-v2.txt": ["Quarterly revenue", "revised", "All regions grew"],
    }
    for name, needles in expected_text.items():
        text = (rendered / name).read_text(encoding="utf-8")
        for needle in needles:
            assert needle in text, (name, needle)
    pngs = sorted(rendered.glob("*.png"))
    assert len(pngs) >= 6
    for font_receipt in sorted(rendered.glob("*.pdffonts.txt")):
        rows = font_receipt.read_text(encoding="utf-8").splitlines()[2:]
        assert rows
        for row in rows:
            columns = row.split()
            assert "Noto" in columns[0]
            assert columns[-5:-2] == ["yes", "yes", "yes"]
    guard = runtime_guard(output / "runtime-guard.tsv")
    assert guard["pack"] == "office-authoring"
    checks = json.loads((Path(__file__).resolve().parents[1] / "pack.lock.json").read_text())[
        "conformance"
    ]["requiredChecks"]
    files = file_receipts(output)
    canonical_json(
        output / "conformance-receipt.json",
        {
            "schema": "ambit.runtime-pack-conformance/v3",
            "packRef": PACK_REF,
            "outcome": "passed",
            "fullImage": True,
            "network": "none",
            "runtime": guard,
            "checks": [{"ref": check, "outcome": "passed"} for check in checks],
            "files": files,
            "limitations": [
                "native-microsoft-office-fidelity-unsupported",
                "windows-office-executor-separate-licensed-profile-required",
                "vba-execution-disabled",
            ],
        },
    )


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("command", choices=("generate", "finalize"))
    parser.add_argument("output", type=Path)
    args = parser.parse_args()
    try:
        if args.command == "generate":
            generate(args.output)
        else:
            finalize(args.output)
    except (AssertionError, OSError, ValueError, zipfile.BadZipFile) as error:
        print(f"office-conformance: {error}", file=sys.stderr)
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
