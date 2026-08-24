from __future__ import annotations

import base64
import csv
import hashlib
import io
import json
import stat
import zipfile
from pathlib import Path
from typing import Any
from xml.etree import ElementTree

from openpyxl import load_workbook
from openpyxl.formula import Tokenizer
from openpyxl.utils.cell import range_boundaries
from PIL import Image
from pptx import Presentation

from process_control import ProcessDeadlineExceeded, ProcessFailure, run_bounded
from render_command import pack_check_names
from render_runner import AdapterFailure


PACK_ROOT = Path("/opt/ambit/runtime-pack/office-authoring")
PATH = f"{PACK_ROOT}/python/bin:/usr/local/bin:/usr/sbin:/usr/bin:/sbin:/bin"
XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
ODS = "application/vnd.oasis.opendocument.spreadsheet"
PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
ODP = "application/vnd.oasis.opendocument.presentation"
MAXIMUM_ARCHIVE_ENTRIES = 20_000
MAXIMUM_ARCHIVE_ENTRY_BYTES = 128 * 1024 * 1024
MAXIMUM_ARCHIVE_TOTAL_BYTES = 2 * 1024 * 1024 * 1024


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as source:
        while chunk := source.read(1024 * 1024):
            digest.update(chunk)
    return "sha256:" + digest.hexdigest()


def _environment(scratch: Path) -> dict[str, str]:
    home = scratch / "home"
    cache = scratch / "cache"
    config = scratch / "config"
    runtime = scratch / "run"
    for directory in (home, cache, config, runtime):
        directory.mkdir(mode=0o700)
    return {
        "HOME": str(home),
        "LANG": "C.UTF-8",
        "LC_ALL": "C.UTF-8",
        "PATH": PATH,
        "PYTHONDONTWRITEBYTECODE": "1",
        "SAL_DISABLE_JAVA": "1",
        "TZ": "UTC",
        "XDG_CACHE_HOME": str(cache),
        "XDG_CONFIG_HOME": str(config),
        "XDG_RUNTIME_DIR": str(runtime),
    }


def _archive(path: Path, required: set[str]) -> dict[str, Any]:
    with zipfile.ZipFile(path) as archive:
        infos = archive.infolist()
        names = [item.filename for item in infos]
        if (
            len(infos) > MAXIMUM_ARCHIVE_ENTRIES
            or len(names) != len(set(names))
            or any(
                item.is_dir() is False
                and (
                    item.file_size > MAXIMUM_ARCHIVE_ENTRY_BYTES
                    or item.filename.startswith("/")
                    or ".." in Path(item.filename).parts
                    or "\\" in item.filename
                    or stat.S_ISLNK(item.external_attr >> 16)
                )
                for item in infos
            )
            or sum(item.file_size for item in infos) > MAXIMUM_ARCHIVE_TOTAL_BYTES
            or not required.issubset(names)
        ):
            raise AdapterFailure(
                "unsafe_office_package",
                "The office package is malformed or exceeds its safe archive bounds.",
                check="office.package_containment",
                observations={"entries": len(infos)},
            )
        return {
            "entryCount": len(infos),
            "uncompressedBytes": sum(item.file_size for item in infos),
            "requiredMembers": sorted(required),
        }


def _run(argv: list[str], *, scratch: Path, deadline: float, environment: dict[str, str]) -> bytes:
    try:
        return run_bounded(
            argv,
            deadline=deadline,
            cwd=scratch,
            environment=environment,
        ).stdout
    except ProcessDeadlineExceeded as error:
        raise AdapterFailure(
            "office_deadline_exceeded",
            "The office renderer exceeded its exact deadline.",
            outcome="blocked",
        ) from error
    except ProcessFailure as error:
        raise AdapterFailure(
            "office_native_tool_failed",
            "An offline office validation tool rejected the artifact.",
            observations={
                "tool": Path(error.result.argv[0]).name,
                "exitCode": error.result.returncode,
                "stderrSha256": "sha256:" + hashlib.sha256(error.result.stderr).hexdigest(),
            },
        ) from error


def _render_pdf(
    source: Path,
    *,
    scratch: Path,
    deadline: float,
    environment: dict[str, str],
) -> tuple[Path, str, list[Path], list[str], Path, Path]:
    rendered = scratch / "rendered"
    profile = scratch / "libreoffice-profile"
    rendered.mkdir()
    profile.mkdir()
    _run(
        [
            "soffice",
            "--headless",
            "--nologo",
            "--nodefault",
            "--nofirststartwizard",
            f"-env:UserInstallation={profile.as_uri()}",
            "--convert-to",
            "pdf",
            "--outdir",
            str(rendered),
            str(source),
        ],
        scratch=scratch,
        deadline=deadline,
        environment=environment,
    )
    pdf = rendered / f"{source.stem}.pdf"
    if not pdf.is_file() or pdf.stat().st_size == 0:
        raise AdapterFailure(
            "office_render_missing",
            "LibreOffice did not produce a complete PDF render.",
            check="office.render_complete",
        )
    text_path = rendered / "text.txt"
    _run(
        ["pdftotext", "-layout", str(pdf), str(text_path)],
        scratch=scratch,
        deadline=deadline,
        environment=environment,
    )
    font_output = _run(
        ["pdffonts", str(pdf)],
        scratch=scratch,
        deadline=deadline,
        environment=environment,
    ).decode("utf-8", errors="strict")
    font_path = rendered / "fonts.txt"
    font_path.write_text(font_output, encoding="utf-8")
    font_path.chmod(0o400)
    prefix = rendered / "page"
    _run(
        ["pdftoppm", "-png", "-r", "96", str(pdf), str(prefix)],
        scratch=scratch,
        deadline=deadline,
        environment=environment,
    )
    pages = sorted(rendered.glob("page-*.png"))
    if not pages:
        raise AdapterFailure(
            "office_page_render_missing",
            "The office renderer produced no page images.",
            check="office.render_complete",
        )
    font_rows = font_output.splitlines()[2:]
    return pdf, text_path.read_text(encoding="utf-8"), pages, font_rows, text_path, font_path


def _font_environment(
    scratch: Path,
    deadline: float,
    environment: dict[str, str],
) -> dict[str, Any]:
    manifest = PACK_ROOT / "locks/fonts/font-files.sha256"
    verified = 0
    for line in manifest.read_text(encoding="utf-8").splitlines():
        expected, separator, relative = line.partition("  ")
        path = Path("/") / relative
        if not separator or not path.is_file() or hashlib.sha256(path.read_bytes()).hexdigest() != expected:
            raise AdapterFailure(
                "font_manifest_drift",
                "The installed font bytes differ from the exact admitted manifest.",
                check="presentation.font_substitution_accounted",
            )
        verified += 1
    actual = _run(
        ["fc-list", "-f", "%{file}\t%{family}\t%{style}\n"],
        scratch=scratch,
        deadline=deadline,
        environment=environment,
    )
    actual_lines = sorted(actual.decode("utf-8", errors="strict").splitlines())
    expected_lines = (
        PACK_ROOT / "locks/fonts/fontconfig-roster.tsv"
    ).read_text(encoding="utf-8").splitlines()
    if actual_lines != expected_lines:
        raise AdapterFailure(
            "fontconfig_roster_drift",
            "The installed fontconfig roster differs from the exact admitted roster.",
            check="presentation.font_substitution_accounted",
        )
    receipt = scratch / "font-environment.json"
    receipt.write_text(
        json.dumps(
            {
                "fontFileCount": verified,
                "fontManifestDigest": _sha256(manifest),
                "fontconfigEntryCount": len(actual_lines),
                "fontconfigRosterDigest": _sha256(
                    PACK_ROOT / "locks/fonts/fontconfig-roster.tsv"
                ),
            },
            indent=2,
            sort_keys=True,
        )
        + "\n",
        encoding="utf-8",
    )
    receipt.chmod(0o400)
    return {
        "fontFileCount": verified,
        "fontconfigEntryCount": len(actual_lines),
        "receipt": receipt,
    }


def _workbook(path: Path, scratch: Path, deadline: float, environment: dict[str, str]) -> dict[str, Any]:
    package = _archive(path, {"[Content_Types].xml", "xl/workbook.xml"})
    workbook = load_workbook(path, data_only=False, read_only=False)
    if not workbook.sheetnames:
        raise AdapterFailure("empty_workbook", "The workbook contains no worksheets.")
    formulas: list[tuple[str, str, str]] = []
    styled_cells = 0
    table_count = 0
    chart_count = 0
    pivot_count = 0
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.style_id:
                    styled_cells += 1
                if isinstance(cell.value, str) and cell.value.startswith("="):
                    Tokenizer(cell.value)
                    formulas.append((sheet.title, cell.coordinate, cell.value))
        for table in sheet.tables.values():
            range_boundaries(table.ref)
            table_count += 1
        chart_count += len(sheet._charts)
        pivot_count += len(sheet._pivots)
    for defined_name in workbook.defined_names.values():
        if defined_name.type == "RANGE":
            for sheet_name, _coordinate in defined_name.destinations:
                if sheet_name not in workbook.sheetnames:
                    raise AdapterFailure(
                        "workbook_reference_invalid",
                        "A workbook defined name targets an absent worksheet.",
                        check="spreadsheet.formula_reference_name_integrity",
                    )
    if not workbook.properties.title or not workbook.properties.subject:
        raise AdapterFailure(
            "workbook_accessibility_metadata_missing",
            "The workbook is missing required title or subject metadata.",
            check="spreadsheet.accessibility_structure",
        )
    for sheet in workbook.worksheets:
        header = [cell.value for cell in next(sheet.iter_rows(min_row=1, max_row=1))]
        if not header or any(value in {None, ""} for value in header):
            raise AdapterFailure(
                "workbook_header_missing",
                "A worksheet has an incomplete first-row header structure.",
                check="spreadsheet.accessibility_structure",
            )

    recalculated = scratch / "recalculated"
    recalc_profile = scratch / "recalc-profile"
    recalculated.mkdir()
    recalc_profile.mkdir()
    _run(
        [
            "soffice",
            "--headless",
            "--nologo",
            "--nodefault",
            "--nofirststartwizard",
            f"-env:UserInstallation={recalc_profile.as_uri()}",
            "--convert-to",
            "xlsx",
            "--outdir",
            str(recalculated),
            str(path),
        ],
        scratch=scratch,
        deadline=deadline,
        environment=environment,
    )
    recalculated_path = recalculated / f"{path.stem}.xlsx"
    if not recalculated_path.is_file():
        raise AdapterFailure(
            "workbook_recalculation_missing",
            "LibreOffice did not emit a recalculated workbook.",
            check="spreadsheet.recalculation_current",
        )
    values = load_workbook(recalculated_path, data_only=True, read_only=True)
    unresolved = [
        f"{sheet}!{coordinate}"
        for sheet, coordinate, _formula in formulas
        if values[sheet][coordinate].value is None
    ]
    if unresolved:
        raise AdapterFailure(
            "workbook_recalculation_incomplete",
            "One or more workbook formulas have no recalculated value.",
            check="spreadsheet.recalculation_current",
            observations={"unresolvedFormulaCount": len(unresolved)},
        )
    return {
        "package": package,
        "sheetCount": len(workbook.sheetnames),
        "formulaCount": len(formulas),
        "styledCellCount": styled_cells,
        "tableCount": table_count,
        "chartCount": chart_count,
        "pivotCount": pivot_count,
        "definedNameCount": len(workbook.defined_names),
        "recalculatedFormulaCount": len(formulas),
    }


def _flat_table(path: Path, media_type: str) -> dict[str, Any]:
    text = path.read_text(encoding="utf-8", errors="strict")
    delimiter = "\t" if media_type == "text/tab-separated-values" else ","
    rows = list(csv.reader(io.StringIO(text, newline=""), delimiter=delimiter, strict=True))
    if not rows or not rows[0] or any(len(row) != len(rows[0]) for row in rows):
        raise AdapterFailure(
            "delimited_structure_invalid",
            "The delimited table is empty or has inconsistent row widths.",
            check="spreadsheet.delimited_structure",
        )
    dangerous = sum(
        value.lstrip().startswith(("=", "+", "-", "@")) for row in rows for value in row
    )
    if dangerous:
        raise AdapterFailure(
            "delimited_formula_injection",
            "The delimited table contains spreadsheet formula-injection cells.",
            check="spreadsheet.formula_injection_policy",
            observations={"dangerousCellCount": dangerous},
        )
    return {
        "rowCount": len(rows),
        "columnCount": len(rows[0]),
        "delimiter": "tab" if delimiter == "\t" else "comma",
        "utf8": True,
        "formulaInjectionCellCount": 0,
    }


def _presentation(path: Path) -> dict[str, Any]:
    package = _archive(
        path,
        {"[Content_Types].xml", "ppt/presentation.xml", "ppt/theme/theme1.xml"},
    )
    presentation = Presentation(path)
    if not presentation.slides:
        raise AdapterFailure("empty_presentation", "The presentation contains no slides.")
    chart_count = 0
    media_count = 0
    note_count = 0
    described_content = 0
    overlap_count = 0
    for slide in presentation.slides:
        if not slide.shapes.title or not slide.shapes.title.text.strip():
            raise AdapterFailure(
                "presentation_title_missing",
                "A presentation slide has no accessible title.",
                check="presentation.accessibility_reading_order",
            )
        notes = slide.notes_slide.notes_text_frame
        if notes is not None and notes.text.strip():
            note_count += 1
        content = []
        described_graphics = 0
        graphics = 0
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > presentation.slide_width or shape.top + shape.height > presentation.slide_height:
                raise AdapterFailure(
                    "presentation_overflow",
                    "A presentation shape exceeds the slide boundary.",
                    check="presentation.layout_overflow_collision",
                )
            if getattr(shape, "has_chart", False):
                chart_count += 1
            if shape.shape_type == 13:
                media_count += 1
            if shape.shape_id != slide.shapes.title.shape_id:
                content.append(shape)
                if not getattr(shape, "has_text_frame", False):
                    graphics += 1
                    descriptions = shape._element.xpath(".//*[local-name()='cNvPr']/@descr")
                    if descriptions and descriptions[0]:
                        described_content += 1
                        described_graphics += 1
        for index, left in enumerate(content):
            for right in content[index + 1 :]:
                if not (
                    left.left + left.width <= right.left
                    or right.left + right.width <= left.left
                    or left.top + left.height <= right.top
                    or right.top + right.height <= left.top
                ):
                    overlap_count += 1
        if graphics and described_graphics != graphics:
            raise AdapterFailure(
                "presentation_alt_text_missing",
                "Presentation content lacks alternative descriptions.",
                check="presentation.accessibility_reading_order",
            )
    return {
        "package": package,
        "slideCount": len(presentation.slides),
        "chartCount": chart_count,
        "mediaCount": media_count,
        "notesCount": note_count,
        "describedContentCount": described_content,
        "collisionSuspectCount": overlap_count,
        "boundaryOverflowCount": 0,
    }


def _open_document(path: Path, media_type: str) -> dict[str, Any]:
    required = {"content.xml", "styles.xml", "META-INF/manifest.xml"}
    package = _archive(path, required)
    with zipfile.ZipFile(path) as archive:
        content = ElementTree.fromstring(archive.read("content.xml"))
        styles = ElementTree.fromstring(archive.read("styles.xml"))
    text_nodes = [value.strip() for value in content.itertext() if value.strip()]
    if not text_nodes:
        raise AdapterFailure("open_document_empty", "The OpenDocument artifact has no text structure.")
    result = {
        "package": package,
        "textNodeCount": len(text_nodes),
        "styleNodeCount": sum(1 for _ in styles.iter()),
    }
    if media_type == ODS:
        result["tableCount"] = sum(1 for node in content.iter() if node.tag.endswith("}table"))
    else:
        result["slideCount"] = sum(1 for node in content.iter() if node.tag.endswith("}page"))
    return result


def _views(
    pages: list[Path],
    text: str,
    request: dict[str, Any],
) -> tuple[list[dict[str, Any]], list[dict[str, str]]]:
    views: list[dict[str, Any]] = []
    limitations: list[dict[str, str]] = []
    aggregate_pixels = 0
    aggregate_bytes = 0
    for page in pages:
        payload = page.read_bytes()
        with Image.open(page) as image:
            width, height = image.size
        pixels = width * height
        if (
            width > 4_096
            or height > 4_096
            or pixels > request["output"]["maximumImagePixels"]
            or aggregate_pixels + pixels > request["output"]["maximumAggregateImagePixels"]
            or len(payload) > 512 * 1024
            or aggregate_bytes + len(payload) > 3 * 1024 * 1024
        ):
            limitations.append(
                {
                    "code": "preview_page_roster_truncated",
                    "severity": "warning",
                    "message": "Additional rendered pages remain in private evidence because the safe public preview reached its image bound.",
                }
            )
            break
        aggregate_pixels += pixels
        aggregate_bytes += len(payload)
        views.append(
            {
                "kind": "image",
                "ordinal": len(views) + 1,
                "label": f"Rendered page {len(views) + 1}",
                "altText": f"Rendered office artifact page {len(views) + 1} of {len(pages)}.",
                "mediaType": "image/png",
                "width": width,
                "height": height,
                "byteLength": len(payload),
                "digest": "sha256:" + hashlib.sha256(payload).hexdigest(),
                "bodyBase64": base64.b64encode(payload).decode(),
            }
        )
    normalized = text.replace("\r", "").replace("\f", "\n")
    encoded = normalized.encode("utf-8")
    if encoded:
        if len(encoded) > 1024 * 1024:
            encoded = encoded[: 1024 * 1024]
            while True:
                try:
                    normalized = encoded.decode("utf-8")
                    break
                except UnicodeDecodeError:
                    encoded = encoded[:-1]
            limitations.append(
                {
                    "code": "preview_text_truncated",
                    "severity": "warning",
                    "message": "The safe public text preview is truncated; complete extraction remains in private evidence.",
                }
            )
        views.append(
            {
                "kind": "text",
                "ordinal": len(views) + 1,
                "label": "Extracted text",
                "mediaType": "text/plain",
                "byteLength": len(encoded),
                "digest": "sha256:" + hashlib.sha256(encoded).hexdigest(),
                "body": normalized,
            }
        )
    return views, sorted(limitations, key=lambda item: item["code"])


def render_validate(
    *,
    request: dict[str, Any],
    source_path: Path,
    scratch: Path,
    deadline: float,
) -> dict[str, Any]:
    media_type = request["source"]["mediaType"]
    environment = _environment(scratch)
    if request["facet"] == "spreadsheet":
        if media_type == XLSX:
            structure = _workbook(source_path, scratch, deadline, environment)
        elif media_type in {"text/csv", "text/tab-separated-values"}:
            structure = _flat_table(source_path, media_type)
        elif media_type == ODS:
            structure = _open_document(source_path, media_type)
        else:
            raise AdapterFailure("unsupported_spreadsheet", "The spreadsheet format is unsupported.")
        title = "Spreadsheet preview"
        summary = "Deterministic structural, recalculation, and page-render validation of the spreadsheet source."
    elif request["facet"] == "presentation":
        if media_type == PPTX:
            structure = _presentation(source_path)
        elif media_type == ODP:
            structure = _open_document(source_path, media_type)
        else:
            raise AdapterFailure("unsupported_presentation", "The presentation format is unsupported.")
        title = "Presentation preview"
        summary = "Deterministic package, accessibility-structure, layout, font, and slide-render validation."
    else:
        raise AdapterFailure("facet_not_owned", "The office pack does not own this artifact facet.")

    font_environment = _font_environment(scratch, deadline, environment)
    pdf, extracted_text, pages, font_rows, text_path, font_path = _render_pdf(
        source_path,
        scratch=scratch,
        deadline=deadline,
        environment=environment,
    )
    fonts = []
    for row in font_rows:
        columns = row.split()
        if len(columns) < 7:
            continue
        fonts.append(columns[0])
        if "Noto" not in columns[0] or columns[-5:-2] != ["yes", "yes", "yes"]:
            raise AdapterFailure(
                "font_substitution_unaccounted",
                "The rendered office artifact uses an unaccounted or unembedded font.",
                check=(
                    "presentation.font_substitution_accounted"
                    if request["facet"] == "presentation"
                    else "spreadsheet.sheet_render_complete"
                ),
                observations={"font": columns[0]},
            )
    common = {
        "sourceDigest": _sha256(source_path),
        "renderedPdfDigest": _sha256(pdf),
        "renderedPageCount": len(pages),
        "fontRoster": sorted(set(fonts)),
        "fontEnvironment": {
            "fontFileCount": font_environment["fontFileCount"],
            "fontconfigEntryCount": font_environment["fontconfigEntryCount"],
            "receiptDigest": _sha256(font_environment["receipt"]),
        },
        **structure,
    }
    supported: dict[str, dict[str, Any]] = {}
    for check in pack_check_names(request):
        supported[check] = {"check": check, **common}
    views, limitations = _views(pages, extracted_text, request)
    facts = [
        {"key": "font_count", "label": "Fonts", "value": str(len(set(fonts)))},
        {"key": "page_count", "label": "Rendered pages", "value": str(len(pages))},
        {"key": "source_bytes", "label": "Source bytes", "value": str(source_path.stat().st_size)},
    ]
    render_check = (
        "presentation.slide_render_complete"
        if request["facet"] == "presentation"
        else "spreadsheet.sheet_render_complete"
    )
    return {
        "title": title,
        "summary": summary,
        "views": views,
        "facts": facts,
        "limitations": limitations,
        "observations": supported,
        "evidenceArtifacts": [
            {
                "check": render_check,
                "mediaType": "application/pdf",
                "name": "rendered.pdf",
                "sourcePath": str(pdf),
            },
            {
                "check": render_check,
                "mediaType": "text/plain",
                "name": "extracted-text.txt",
                "sourcePath": str(text_path),
            },
            {
                "check": "presentation.font_substitution_accounted"
                if request["facet"] == "presentation"
                else render_check,
                "mediaType": "application/json",
                "name": "font-environment.json",
                "sourcePath": str(font_environment["receipt"]),
            },
            {
                "check": "presentation.font_substitution_accounted"
                if request["facet"] == "presentation"
                else render_check,
                "mediaType": "text/plain",
                "name": "rendered-fonts.txt",
                "sourcePath": str(font_path),
            },
        ],
    }
